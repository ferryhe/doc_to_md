"""Tests for PPTX and XLSX text extraction functions."""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest

pytest.importorskip("pptx", reason="python-pptx not installed")
pytest.importorskip("openpyxl", reason="openpyxl not installed")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_pptx(tmp_path: Path, slides: list[dict]) -> Path:
    """Build a minimal .pptx file with the given slide content."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content
    for slide_data in slides:
        slide = prs.slides.add_slide(layout)
        if "title" in slide_data:
            slide.shapes.title.text = slide_data["title"]
        if "body" in slide_data:
            slide.placeholders[1].text = slide_data["body"]
    buf = BytesIO()
    prs.save(buf)
    pptx_path = tmp_path / "test.pptx"
    pptx_path.write_bytes(buf.getvalue())
    return pptx_path


def _make_xlsx(tmp_path: Path, sheets: dict[str, list[list]]) -> Path:
    """Build a minimal .xlsx file with the given sheet data."""
    import openpyxl

    wb = openpyxl.Workbook()
    first = True
    for sheet_name, rows in sheets.items():
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(sheet_name)
        for row in rows:
            ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    xlsx_path = tmp_path / "test.xlsx"
    xlsx_path.write_bytes(buf.getvalue())
    return xlsx_path


# ---------------------------------------------------------------------------
# _extract_pptx tests
# ---------------------------------------------------------------------------


def test_extract_pptx_returns_markdown(tmp_path: Path) -> None:
    from doc_to_md.pipeline.text_extraction import _extract_pptx

    pptx_path = _make_pptx(tmp_path, [{"title": "My Slide", "body": "Slide content here"}])
    result = _extract_pptx(pptx_path)
    assert "My Slide" in result
    assert "Slide content here" in result


def test_extract_pptx_multiple_slides(tmp_path: Path) -> None:
    from doc_to_md.pipeline.text_extraction import _extract_pptx

    pptx_path = _make_pptx(
        tmp_path,
        [
            {"title": "Slide One", "body": "First"},
            {"title": "Slide Two", "body": "Second"},
        ],
    )
    result = _extract_pptx(pptx_path)
    assert "Slide 1" in result
    assert "Slide 2" in result
    assert "Slide One" in result
    assert "Slide Two" in result


def test_extract_pptx_empty_presentation(tmp_path: Path) -> None:
    from pptx import Presentation

    from doc_to_md.pipeline.text_extraction import _extract_pptx

    prs = Presentation()
    buf = BytesIO()
    prs.save(buf)
    pptx_path = tmp_path / "empty.pptx"
    pptx_path.write_bytes(buf.getvalue())

    result = _extract_pptx(pptx_path)
    assert "No textual content" in result


def test_extract_pptx_missing_dependency(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    import sys
    import doc_to_md.pipeline.text_extraction as mod

    monkeypatch.setitem(sys.modules, "pptx", None)
    pptx_path = tmp_path / "fake.pptx"
    pptx_path.write_bytes(b"PK\x03\x04")

    with pytest.raises(RuntimeError, match="python-pptx"):
        mod._extract_pptx(pptx_path)


# ---------------------------------------------------------------------------
# _extract_xlsx tests
# ---------------------------------------------------------------------------


def test_extract_xlsx_returns_table(tmp_path: Path) -> None:
    from doc_to_md.pipeline.text_extraction import _extract_xlsx

    xlsx_path = _make_xlsx(
        tmp_path,
        {"Results": [["Name", "Score"], ["Alice", "95"], ["Bob", "88"]]},
    )
    result = _extract_xlsx(xlsx_path)
    assert "Name" in result
    assert "Alice" in result
    assert "95" in result


def test_extract_xlsx_multiple_sheets(tmp_path: Path) -> None:
    from doc_to_md.pipeline.text_extraction import _extract_xlsx

    xlsx_path = _make_xlsx(
        tmp_path,
        {
            "Q1": [["Month", "Revenue"], ["Jan", "1000"]],
            "Q2": [["Month", "Revenue"], ["Apr", "1200"]],
        },
    )
    result = _extract_xlsx(xlsx_path)
    assert "## Q1" in result
    assert "## Q2" in result
    assert "Jan" in result
    assert "Apr" in result


def test_extract_xlsx_empty_workbook(tmp_path: Path) -> None:
    import openpyxl

    from doc_to_md.pipeline.text_extraction import _extract_xlsx

    wb = openpyxl.Workbook()
    buf = BytesIO()
    wb.save(buf)
    xlsx_path = tmp_path / "empty.xlsx"
    xlsx_path.write_bytes(buf.getvalue())

    result = _extract_xlsx(xlsx_path)
    assert "No data found" in result


def test_extract_xlsx_missing_dependency(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    import sys
    import doc_to_md.pipeline.text_extraction as mod

    monkeypatch.setitem(sys.modules, "openpyxl", None)
    xlsx_path = tmp_path / "fake.xlsx"
    xlsx_path.write_bytes(b"PK\x03\x04")

    with pytest.raises(RuntimeError, match="openpyxl"):
        mod._extract_xlsx(xlsx_path)


# ---------------------------------------------------------------------------
# Integration: LocalEngine converts pptx / xlsx end-to-end
# ---------------------------------------------------------------------------


def test_local_engine_converts_pptx(tmp_path: Path) -> None:
    from doc_to_md.engines.local import LocalEngine

    pptx_path = _make_pptx(tmp_path, [{"title": "Engine Test", "body": "PPTX via LocalEngine"}])
    response = LocalEngine().convert(pptx_path)
    assert "Engine Test" in response.markdown or "PPTX via LocalEngine" in response.markdown


def test_local_engine_converts_xlsx(tmp_path: Path) -> None:
    from doc_to_md.engines.local import LocalEngine

    xlsx_path = _make_xlsx(tmp_path, {"Data": [["Item", "Qty"], ["Widget", "10"]]})
    response = LocalEngine().convert(xlsx_path)
    assert "Item" in response.markdown or "Widget" in response.markdown


def test_auto_engine_converts_pptx(tmp_path: Path) -> None:
    from doc_to_md.engines.auto import AutoEngine

    pptx_path = _make_pptx(tmp_path, [{"title": "Auto PPTX", "body": "Routed via auto"}])
    response = AutoEngine().convert(pptx_path)
    assert "Auto PPTX" in response.markdown or "Routed via auto" in response.markdown


def test_auto_engine_converts_xlsx(tmp_path: Path) -> None:
    from doc_to_md.engines.auto import AutoEngine

    xlsx_path = _make_xlsx(tmp_path, {"Sheet1": [["Col1", "Col2"], ["A", "B"]]})
    response = AutoEngine().convert(xlsx_path)
    assert "Col1" in response.markdown or "A" in response.markdown
